"""LLM 담당(전승우) 발표용 PPT — B(임베딩) 슬라이드 스타일을 그대로 재현.

실행: python make_ppt_LLM.py
출력: Project_team1/A_LLM_tuning_py/전승우_LLM_발표자료.pptx
구성: 실험설계 / 결과①정확도 / 결과②응답시간·안정성 / 핵심코드① / 핵심코드② / 인사이트  (총 6장)
"""

import os

import matplotlib
matplotlib.use("Agg")
matplotlib.rcParams["font.family"] = "Malgun Gothic"
matplotlib.rcParams["axes.unicode_minus"] = False
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

HERE = os.path.dirname(__file__)
CHART1_PATH = os.path.join(HERE, "_chart_f1.png")
CHART2A_PATH = os.path.join(HERE, "_chart_time.png")
CHART2B_PATH = os.path.join(HERE, "_chart_std.png")
OUT_PATH = r"C:\Users\KDT21\Desktop\전승우_LLM_발표자료.pptx"

# ---- 실측 데이터 (results_A_llm_comparison_all.csv 집계) ----
MODELS = [
    {"name": "Phi-4-mini-instruct", "full": "microsoft/Phi-4-mini-instruct", "f1": 0.6462, "time": 12.70, "std": 0.0703},
    {"name": "Qwen2.5-7B-Instruct\n(baseline)", "full": "Qwen/Qwen2.5-7B-Instruct", "f1": 0.6039, "time": 15.96, "std": 0.0708},
    {"name": "Qwen2.5-3B-Instruct", "full": "Qwen/Qwen2.5-3B-Instruct", "f1": 0.5621, "time": 11.36, "std": 0.1173},
]

# ---- 색상 팔레트 (B 슬라이드와 동일) ----
DARK_TAG = RGBColor(0x33, 0x33, 0x33)
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ORANGE = RGBColor(0xF2, 0x99, 0x4A)
ORANGE_TEXT = RGBColor(0xE0, 0x7A, 0x1E)
LAVENDER = RGBColor(0xF0, 0xF1, 0xF8)
CARD_BG = RGBColor(0xF6, 0xF6, 0xF8)
GRAY_TEXT = RGBColor(0x8A, 0x8A, 0x8A)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x39)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_FG = RGBColor(0xD4, 0xD4, 0xD4)

# ============================================================
# 차트 이미지 생성
# ============================================================
names = [m["name"] for m in MODELS]
f1s = [m["f1"] for m in MODELS]
times = [m["time"] for m in MODELS]
stds = [m["std"] for m in MODELS]
bar_colors_f1 = ["#F2994A" if i == 0 else "#1B2A4A" for i in range(len(MODELS))]

# 결과① 정확도 — 가로 막대
fig, ax = plt.subplots(figsize=(6.6, 3.6), dpi=200)
y_pos = range(len(names))
bars = ax.barh(list(y_pos)[::-1], f1s, color=bar_colors_f1, height=0.55)
for i, (b, v) in enumerate(zip(bars, f1s)):
    ax.text(v + 0.012, b.get_y() + b.get_height() / 2, f"{v:.4f}", va="center", fontsize=12, fontweight="bold")
ax.set_yticks(list(y_pos)[::-1])
ax.set_yticklabels(names, fontsize=11)
ax.set_xlim(0, 0.78)
ax.spines[["top", "right", "left"]].set_visible(False)
ax.xaxis.grid(True, color="#e2e2e2", linewidth=0.8)
ax.set_axisbelow(True)
ax.tick_params(axis="x", labelsize=9)
plt.tight_layout()
plt.savefig(CHART1_PATH, transparent=True)
plt.close()

# 결과② 응답시간 — 세로 막대(전부 오렌지)
fig, ax = plt.subplots(figsize=(4.6, 3.9), dpi=200)
bars = ax.bar(range(len(names)), times, color="#F2994A", width=0.55)
for b, v in zip(bars, times):
    ax.text(b.get_x() + b.get_width() / 2, v + 0.3, f"{v:.1f}", ha="center", fontsize=11, fontweight="bold")
ax.set_xticks(range(len(names)))
ax.set_xticklabels(names, fontsize=9)
ax.set_ylim(0, max(times) * 1.25)
ax.spines[["top", "right", "left"]].set_visible(False)
ax.yaxis.grid(True, color="#e2e2e2", linewidth=0.8)
ax.set_axisbelow(True)
plt.tight_layout()
plt.savefig(CHART2A_PATH, transparent=True)
plt.close()

# 결과② F1 표준편차 — 세로 막대(최고 성능 모델 오렌지 강조)
fig, ax = plt.subplots(figsize=(4.6, 3.9), dpi=200)
bars = ax.bar(range(len(names)), stds, color=bar_colors_f1, width=0.55)
for b, v in zip(bars, stds):
    ax.text(b.get_x() + b.get_width() / 2, v + 0.003, f"{v:.3f}", ha="center", fontsize=11, fontweight="bold")
ax.set_xticks(range(len(names)))
ax.set_xticklabels(names, fontsize=9)
ax.set_ylim(0, max(stds) * 1.3)
ax.spines[["top", "right", "left"]].set_visible(False)
ax.yaxis.grid(True, color="#e2e2e2", linewidth=0.8)
ax.set_axisbelow(True)
plt.tight_layout()
plt.savefig(CHART2B_PATH, transparent=True)
plt.close()

# ============================================================
# pptx 뼈대
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_kdt_tag(slide):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.55), Inches(0.15), Inches(2.2), Inches(0.42))
    tag.adjustments[0] = 0.5
    tag.fill.solid()
    tag.fill.fore_color.rgb = DARK_TAG
    tag.line.fill.background()
    tag.shadow.inherit = False
    p = tag.text_frame.paragraphs[0]
    p.text = "KDT-14기"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE


def add_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.88), Inches(0.06), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = slide.shapes.add_textbox(Inches(0.75), Inches(0.72), Inches(11.0), Inches(0.55))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = DARK_TEXT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.78), Inches(1.28), Inches(11.5), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12.5)
        sp.font.color.rgb = GRAY_TEXT


def add_pagenum(slide, n):
    box = slide.shapes.add_textbox(Inches(12.7), Inches(7.12), Inches(0.5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = str(n)
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT


def add_code_box(slide, code_lines, left, top, width, height, title=None):
    if title:
        tbox = slide.shapes.add_textbox(left, top - Inches(0.32), width, Inches(0.3))
        tp = tbox.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(12)
        tp.font.bold = True
        tp.font.color.rgb = GRAY_TEXT

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.03
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0x3C, 0x3C, 0x3C)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.12)
    tf.margin_right = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.font.color.rgb = CODE_FG
    return box


def add_bullets(slide, items, left, top, width, height, size=13, color=DARK_TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


# ============================================================
# 슬라이드 1 — 실험 설계
# ============================================================
s1 = add_slide()
add_kdt_tag(s1)
add_header(s1, "실험 설계", "LLM 모델 3종을 동일한 임베딩·청킹·검색 조건에서 비교합니다")

navy_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(4.6), Inches(5.0))
navy_box.adjustments[0] = 0.03
navy_box.fill.solid()
navy_box.fill.fore_color.rgb = NAVY
navy_box.line.fill.background()
navy_box.shadow.inherit = False
tf = navy_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.3)

def navy_kv(tf, label, value, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = label
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.space_before = Pt(0 if first else 14)
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0xD8, 0xDE, 0xEB)

navy_kv(tf, "공통 고정 조건", "", first=True)
navy_kv(tf, "임베딩", "jhgan/ko-sroberta-multitask")
navy_kv(tf, "청킹", "마크다운 기반, 500자 / overlap 50")
navy_kv(tf, "검색 방식", "similarity, top-k = 3")
navy_kv(tf, "평가 데이터", "자동차관리법 질의 10문항")
navy_kv(tf, "평가 지표", "BERTScore F1 (klue/bert-base)")

tbl_title = s1.shapes.add_textbox(Inches(5.5), Inches(1.7), Inches(7.2), Inches(0.35))
tp = tbl_title.text_frame.paragraphs[0]
tp.text = "비교 대상 LLM 모델 (3종 유효 결과)"
tp.font.bold = True
tp.font.size = Pt(15)
tp.font.color.rgb = DARK_TEXT

rows = [
    ("Phi-4-mini-instruct", "Microsoft, 3.8B", "0.6462", True),
    ("Qwen2.5-7B-Instruct (baseline)", "Alibaba, 7B", "0.6039", False),
    ("Qwen2.5-3B-Instruct", "Alibaba, 3B", "0.5621", False),
]
gtable = s1.shapes.add_table(len(rows) + 1, 3, Inches(5.5), Inches(2.15), Inches(7.2), Inches(1.7)).table
gtable.columns[0].width = Inches(3.6)
gtable.columns[1].width = Inches(2.0)
gtable.columns[2].width = Inches(1.6)
headers = ["모델", "유형", "F1"]
for c, h in enumerate(headers):
    cell = gtable.cell(0, c)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    pp = cell.text_frame.paragraphs[0]
    pp.font.bold = True
    pp.font.size = Pt(12)
    pp.font.color.rgb = WHITE
for r, (name, typ, f1, best) in enumerate(rows, start=1):
    for c, val in enumerate([name, typ, f1]):
        cell = gtable.cell(r, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if r % 2 else LAVENDER
        pp = cell.text_frame.paragraphs[0]
        pp.font.size = Pt(12)
        pp.font.bold = best
        pp.font.color.rgb = ORANGE_TEXT if (best and c == 2) else DARK_TEXT

add_bullets(s1, [
    "실행 환경: 로컬 GPU 8GB VRAM, 4bit 양자화(bitsandbytes)",
    "2차 확장 후보 3종(Qwen2.5-14B, Bllossom-8B, EXAONE-3.5-2.4B)은",
    "  8GB VRAM 한계·아키텍처 호환성 문제로 로드/실행 실패",
], Inches(5.5), Inches(4.1), Inches(7.2), Inches(2.2), size=13)
add_pagenum(s1, 1)

# ============================================================
# 슬라이드 2 — 결과① 정확도
# ============================================================
s2 = add_slide()
add_kdt_tag(s2)
add_header(s2, "결과 ① 정확도 (BERTScore F1)", "동일 질의 10개에 대한 답변 품질 — 값이 높을수록 정답에 가까움")

s2.shapes.add_picture(CHART1_PATH, Inches(0.5), Inches(1.8), width=Inches(7.4))

tag = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.85), Inches(1.6), Inches(0.4))
tag.adjustments[0] = 0.5
tag.fill.solid()
tag.fill.fore_color.rgb = ORANGE
tag.line.fill.background()
tag.shadow.inherit = False
tp = tag.text_frame.paragraphs[0]
tp.text = "최고 성능"
tp.font.size = Pt(12)
tp.font.bold = True
tp.font.color.rgb = WHITE
tp.alignment = PP_ALIGN.CENTER

num_box = s2.shapes.add_textbox(Inches(8.3), Inches(2.35), Inches(4.3), Inches(0.9))
p = num_box.text_frame.paragraphs[0]
p.text = "0.6462"
p.font.size = Pt(42)
p.font.bold = True
p.font.color.rgb = ORANGE_TEXT
p2 = num_box.text_frame.add_paragraph()
p2.text = "Phi-4-mini-instruct (baseline 0.6039 대비 +7.0%)"
p2.font.size = Pt(12)
p2.font.color.rgb = DARK_TEXT

top3_box = s2.shapes.add_textbox(Inches(8.3), Inches(3.5), Inches(4.3), Inches(1.6))
tf = top3_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Top 3"
p.font.bold = True
p.font.size = Pt(15)
p.font.color.rgb = DARK_TEXT
for i, m in enumerate(MODELS, start=1):
    pp = tf.add_paragraph()
    pp.text = f"{i}. {m['name'].splitlines()[0]}  {m['f1']:.4f}"
    pp.font.size = Pt(13)
    pp.font.color.rgb = DARK_TEXT

insight_card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(5.15), Inches(4.3), Inches(1.4))
insight_card.adjustments[0] = 0.08
insight_card.fill.solid()
insight_card.fill.fore_color.rgb = LAVENDER
insight_card.line.fill.background()
insight_card.shadow.inherit = False
itf = insight_card.text_frame
itf.word_wrap = True
itf.margin_left = Inches(0.2)
itf.vertical_anchor = MSO_ANCHOR.MIDDLE
ip = itf.paragraphs[0]
ip.text = "1차: 7B가 3B를 유의미하게 앞섬 → 2차 확장에서 소형(3.8B) Phi-4-mini가 baseline 7B를 역전"
ip.font.size = Pt(12.5)
ip.font.color.rgb = DARK_TEXT
add_pagenum(s2, 2)

# ============================================================
# 슬라이드 3 — 결과② 응답시간·안정성
# ============================================================
s3 = add_slide()
add_kdt_tag(s3)
add_header(s3, "결과 ② 응답시간 · 답변 안정성", "정확도 외에 속도와 질문별 편차(표준편차)도 함께 확인합니다")

t1 = s3.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(5.5), Inches(0.35))
t1.text_frame.paragraphs[0].text = "평균 응답시간 (초)"
t1.text_frame.paragraphs[0].font.bold = True
t1.text_frame.paragraphs[0].font.size = Pt(13)
s3.shapes.add_picture(CHART2A_PATH, Inches(0.6), Inches(2.15), width=Inches(5.7))

t2 = s3.shapes.add_textbox(Inches(6.9), Inches(1.75), Inches(5.8), Inches(0.35))
t2.text_frame.paragraphs[0].text = "질문별 F1 표준편차 (낮을수록 안정적)"
t2.text_frame.paragraphs[0].font.bold = True
t2.text_frame.paragraphs[0].font.size = Pt(13)
s3.shapes.add_picture(CHART2B_PATH, Inches(6.9), Inches(2.15), width=Inches(5.7))

add_bullets(s3, [
    "Phi-4-mini는 baseline(7B)보다 3.3초 더 빠르면서도 F1은 더 높음 — 속도·정확도 동시 개선",
    "Qwen2.5-3B는 가장 빠르지만(11.4초) 표준편차 0.117로 가장 불안정 — 질문에 따라 품질 기복이 큼",
], Inches(0.7), Inches(6.4), Inches(12.0), Inches(1.0), size=12.5)
add_pagenum(s3, 3)

# ============================================================
# 슬라이드 4 — 카테고리별 강약점
# ============================================================
s4 = add_slide()
add_kdt_tag(s4)
add_header(s4, "카테고리별 강약점", "10개 질문 유형별로 쪼개보면, Phi-4-mini가 전 영역에서 이긴 것은 아니었음")

cat_rows = [
    ("말소 등록 (폐차 등)", "0.369", "0.550", "0.637", 2),
    ("명의이전(이전등록)", "0.516", "0.532", "0.661", 2),
    ("불법 튜닝 및 처벌", "0.428", "0.460", "0.505", 2),
    ("이륜자동차(오토바이) 관리", "0.650", "0.626", "0.590", 0),
    ("자동차 등록 번호판", "0.582", "0.630", "0.644", 2),
    ("자율주행차 및 임시운행허가", "0.706", "0.700", "0.750", 2),
    ("정기검사 주기 및 과태료", "0.450", "0.640", "0.690", 2),
    ("제작결함 시정(리콜)", "0.669", "0.670", "0.648", 1),
    ("중고차 성능·상태 점검", "0.575", "0.613", "0.607", 1),
    ("튜닝 관련", "0.678", "0.621", "0.729", 2),
]
headers = ["질문 카테고리", "Qwen2.5-3B", "Qwen2.5-7B", "Phi-4-mini"]
gtable = s4.shapes.add_table(len(cat_rows) + 1, 4, Inches(0.6), Inches(1.8), Inches(12.1), Inches(4.55)).table
gtable.columns[0].width = Inches(4.9)
gtable.columns[1].width = Inches(2.4)
gtable.columns[2].width = Inches(2.4)
gtable.columns[3].width = Inches(2.4)
for c, h in enumerate(headers):
    cell = gtable.cell(0, c)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    pp = cell.text_frame.paragraphs[0]
    pp.font.bold = True
    pp.font.size = Pt(12)
    pp.font.color.rgb = WHITE
    pp.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
for r, (cat, v3, v7, vphi, win) in enumerate(cat_rows, start=1):
    vals = [cat, v3, v7, vphi]
    for c, val in enumerate(vals):
        cell = gtable.cell(r, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if r % 2 else LAVENDER
        pp = cell.text_frame.paragraphs[0]
        pp.font.size = Pt(11.5)
        pp.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        is_winner_cell = (c - 1) == win
        pp.font.bold = is_winner_cell
        pp.font.color.rgb = ORANGE_TEXT if is_winner_cell else DARK_TEXT
add_pagenum(s4, 4)

# ============================================================
# 슬라이드 5 — 단독 승자가 팀 결합 시엔 최종 승자가 아니었음 (최종선택 패턴)
# ============================================================
s5 = add_slide()
add_kdt_tag(s5)
add_header(s5, "단독 승자가 팀 결합에선 기각됨")

sub = s5.shapes.add_textbox(Inches(0.78), Inches(1.32), Inches(11.5), Inches(0.4))
sp = sub.text_frame.paragraphs[0]
sp.text = "Greedy Search STEP1 — LLM 교체 시도"
sp.font.size = Pt(15)
sp.font.color.rgb = GRAY_TEXT

stat_card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.95), Inches(5.4), Inches(4.5))
stat_card.adjustments[0] = 0.04
stat_card.fill.solid()
stat_card.fill.fore_color.rgb = NAVY
stat_card.line.fill.background()
stat_card.shadow.inherit = False
stf = stat_card.text_frame
stf.word_wrap = True
stf.margin_top = Inches(0.35)

def stat_line(tf, label, value, value_color, first=False, big=True):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xB9, 0xC1, 0xD6)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0 if first else 22)
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.bold = True
    p2.font.size = Pt(34 if big else 20)
    p2.font.color.rgb = value_color
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)

stat_line(stf, "KoE5 + Phi-4-mini (A 단독 승자로 교체)", "0.6618", ORANGE, first=True)
stat_line(stf, "KoE5 + Qwen2.5-7B (baseline 유지)", "0.6715", WHITE)
stat_line(stf, "판정", "LLM 교체 기각", RGBColor(0xE8, 0x5A, 0x5A), big=False)

reason_card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(1.95), Inches(6.4), Inches(4.5))
reason_card.adjustments[0] = 0.04
reason_card.fill.solid()
reason_card.fill.fore_color.rgb = CARD_BG
reason_card.line.fill.background()
reason_card.shadow.inherit = False
rtf = reason_card.text_frame
rtf.word_wrap = True
rtf.margin_left = Inches(0.3)
rtf.margin_top = Inches(0.3)
rp = rtf.paragraphs[0]
rp.text = "왜 이런 일이 벌어졌나"
rp.font.bold = True
rp.font.size = Pt(16)
rp.font.color.rgb = DARK_TEXT

reasons = [
    "A(LLM) 축 비교에서 Phi-4-mini가 이긴 건 임베딩이 jhgan으로 고정된 조건에서였음 (F1 0.6462 vs baseline 0.6039)",
    "팀이 임베딩을 B 단독 승자인 KoE5로 바꾼 뒤 같은 비교를 다시 하니, 오히려 Qwen2.5-7B 쪽이 더 높게 나와 순위가 뒤집힘",
    "즉 \"이 LLM이 좋다\"는 결론은 특정 임베딩과 짝지어졌을 때만 유효 — 축 간 상호작용이 실제로 존재함을 데이터로 확인",
    "그래서 최종 파이프라인의 LLM은 A 단독 승자가 아닌 baseline Qwen2.5-7B로 유지됨",
]
for r in reasons:
    p = rtf.add_paragraph()
    run = p.add_run()
    run.text = "✓ "
    run.font.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = ORANGE_TEXT
    run2 = p.add_run()
    run2.text = r
    run2.font.size = Pt(13)
    run2.font.color.rgb = DARK_TEXT
    p.space_before = Pt(16)

add_pagenum(s5, 5)

# ============================================================
# 슬라이드 6 — 인사이트
# ============================================================
s6 = add_slide()
add_kdt_tag(s6)
add_header(s6, "인사이트", "실험 결과에서 확인한 3가지 시사점")

insights = [
    ("01", "작은 모델이 큰 모델을 이길 수 있다",
     "3.8B인 Phi-4-mini가 7B baseline을 F1 기준 +7.0% 앞섬. 파라미터 수보다 학습 데이터 품질 전략이 이 도메인(법령 QA)에서는 더 결정적이었음."),
    ("02", "속도와 정확도가 항상 트레이드오프는 아니었다",
     "Phi-4-mini는 baseline보다 3.3초 더 빠르면서 F1도 더 높음 — 무조건 큰 모델을 쓸 이유가 없다는 것을 실측으로 확인."),
    ("03", "대형 모델은 로컬 8GB GPU의 한계를 넘기 쉽다",
     "14B/8B급 3종이 로드·실행 단계에서 실패. try/except 방어 코드로 실패가 전체 실험을 중단시키지 않도록 설계해 나머지 결과를 안전하게 확보."),
]

top = Inches(1.85)
for num, title, body in insights:
    circ = s6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.12), Inches(0.7), Inches(0.7))
    circ.fill.solid()
    circ.fill.fore_color.rgb = NAVY
    circ.line.fill.background()
    circ.shadow.inherit = False
    cp = circ.text_frame.paragraphs[0]
    cp.text = num
    cp.font.size = Pt(18)
    cp.font.bold = True
    cp.font.color.rgb = ORANGE
    cp.alignment = PP_ALIGN.CENTER
    circ.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.65), top, Inches(11.0), Inches(1.35))
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = LAVENDER
    card.line.fill.background()
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(14.5)
    p.font.color.rgb = DARK_TEXT
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = GRAY_TEXT
    top = top + Inches(1.55)

add_pagenum(s6, 6)

prs.save(OUT_PATH)
for f in (CHART1_PATH, CHART2A_PATH, CHART2B_PATH):
    os.remove(f)
print(f"저장 완료: {OUT_PATH}")
